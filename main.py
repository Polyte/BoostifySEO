from fastapi import FastAPI, HTTPException, BackgroundTasks
from pydantic import BaseModel
import requests
from bs4 import BeautifulSoup, Tag
import os
from dotenv import load_dotenv
from fastapi.responses import FileResponse, StreamingResponse
import pathlib
import re
from urllib.parse import urljoin, urlparse
import threading
import time
import xml.etree.ElementTree as ET
from fastapi.staticfiles import StaticFiles
from typing import List, Dict, Optional
import json
from datetime import datetime, timedelta
import sqlite3
import hashlib
# New imports for enhanced export
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
try:
    import schedule
    SCHEDULE_AVAILABLE = True
except ImportError:
    SCHEDULE_AVAILABLE = False
import tempfile
import zipfile
from io import BytesIO
import asyncio
import httpx
from collections import defaultdict
import csv
import io
from apscheduler.schedulers.background import BackgroundScheduler
from apscheduler.jobstores.sqlalchemy import SQLAlchemyJobStore
from apscheduler.triggers.cron import CronTrigger
from contextlib import asynccontextmanager

# Load environment variables
load_dotenv()
DEEPSEEK_API_KEY = os.getenv("DEEPSEEK_API_KEY")
print("[DEBUG] DEEPSEEK_API_KEY:", repr(DEEPSEEK_API_KEY))
DEEPSEEK_API_URL = "https://api.deepseek.com/v1/chat/completions"

@asynccontextmanager
async def lifespan(app: FastAPI):
    # Startup code
    if not scheduler.running:
        scheduler.start()
    yield
    # (Optional) Shutdown code can go here

app = FastAPI(lifespan=lifespan)

# Serve all static files (JS, CSS, images, etc.) from the root directory
app.mount("/static", StaticFiles(directory="."), name="static")

# Initialize SQLite database for analytics
def init_database():
    conn = sqlite3.connect('boostify_analytics.db')
    cursor = conn.cursor()
    
    # Create analytics table
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS seo_audits (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            url TEXT NOT NULL,
            audit_data TEXT NOT NULL,
            seo_score INTEGER,
            content_score INTEGER,
            technical_score INTEGER,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            user_id TEXT DEFAULT 'anonymous'
        )
    ''')
    
    # Create competitor analysis table
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS competitor_analyses (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            target_url TEXT NOT NULL,
            competitor_urls TEXT NOT NULL,
            analysis_data TEXT NOT NULL,
            rank_position INTEGER,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            user_id TEXT DEFAULT 'anonymous'
        )
    ''')
    
    # Create trends table
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS seo_trends (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            url TEXT NOT NULL,
            date DATE NOT NULL,
            seo_score INTEGER,
            content_score INTEGER,
            technical_score INTEGER,
            keywords_count INTEGER,
            word_count INTEGER,
            internal_links INTEGER,
            external_links INTEGER,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    
    # Create scheduled reports table
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS scheduled_reports (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            report_id TEXT NOT NULL,
            url TEXT NOT NULL,
            format TEXT NOT NULL,
            report_data BLOB,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    
    conn.commit()
    conn.close()

# Initialize database on startup
init_database()

# Initialize APScheduler
scheduler = BackgroundScheduler(jobstores={
    'default': SQLAlchemyJobStore(url='sqlite:///boostify_analytics.db')
})

# Serve index.html at root
@app.get("/")
def read_index():
    return FileResponse("index.html")

# Serve all HTML pages
@app.get("/single-audit.html")
def read_single_audit():
    return FileResponse("single-audit.html")

@app.get("/full-scan.html")
def read_full_scan():
    return FileResponse("full-scan.html")

@app.get("/competitor-analysis.html")
def read_competitor_analysis():
    return FileResponse("competitor-analysis.html")

@app.get("/analytics.html")
def read_analytics():
    return FileResponse("analytics.html")

@app.get("/export.html")
def read_export():
    return FileResponse("export.html")

# Serve results.html at /results.html
@app.get("/results.html")
def read_results():
    return FileResponse("results.html")

# Serve favicon if it exists
@app.get("/favicon.ico")
def favicon():
    if os.path.exists("favicon.ico"):
        return FileResponse("favicon.ico")
    return "", 204

class SEOAuditRequest(BaseModel):
    url: str

class SEOAuditResponse(BaseModel):
    keywords: list[str]
    meta_description: str
    title_tag: str | None = None
    h1: str | None = None
    word_count: int | None = None
    readability: str | None = None
    missing_alt_text: str | None = None
    suggestions: list[str] | None = None
    # Advanced metrics
    internal_links: int | None = None
    external_links: int | None = None
    image_count: int | None = None
    mobile_friendly: bool | None = None
    page_speed_score: int | None = None
    structured_data: bool | None = None
    canonical_tag: bool | None = None
    robots_meta: bool | None = None
    social_meta: bool | None = None
    # Broken links and images
    broken_images: list[str] | None = None
    broken_links: list[str] | None = None
    # Images with missing alt text
    images_missing_alt: list[dict] | None = None
    # All internal links and images
    all_internal_links: list[dict] | None = None
    all_images: list[dict] | None = None
    # Competitor analysis metrics
    seo_score: int | None = None
    content_quality_score: int | None = None
    technical_score: int | None = None
    # New fields for focus keyphrase and synonyms
    focus_keyphrase: str | None = None
    keyphrase_synonyms: list[str] | None = None
    # Advanced keyphrase analysis
    keyphrase_density: float | None = None
    keyphrase_placement: dict | None = None
    synonym_usage: dict | None = None
    lsi_keywords: list[str] | None = None

class CompetitorAnalysisRequest(BaseModel):
    target_url: str
    competitor_urls: List[str]

class AnalyticsRequest(BaseModel):
    url: str
    date_range: Optional[str] = "30d"  # 7d, 30d, 90d, 1y

class ExportRequest(BaseModel):
    url: str
    format: str  # "pdf", "pptx", "html", "json"
    template: Optional[str] = "default"  # "default", "executive", "detailed"
    include_charts: bool = True
    include_recommendations: bool = True
    include_competitor_analysis: bool = False

class ScheduledReportRequest(BaseModel):
    url: str
    frequency: str  # "daily", "weekly", "monthly"
    format: str  # "pdf", "pptx", "html"
    email: Optional[str] = None
    template: str = "default"

class CompetitorAnalysisResponse(BaseModel):
    target_analysis: SEOAuditResponse
    competitor_analyses: List[SEOAuditResponse]
    comparison_metrics: Dict
    recommendations: List[str]

# In-memory storage for scan jobs (for demo; use a DB for production)
scan_jobs = {}
competitor_jobs = {}

# Helper: Save audit to database
def save_audit_to_db(url: str, audit_data: dict, user_id: str = "anonymous"):
    conn = sqlite3.connect('boostify_analytics.db')
    cursor = conn.cursor()
    
    seo_score = calculate_seo_score_from_dict(audit_data)
    content_score = calculate_content_quality_score_from_dict(audit_data)
    technical_score = calculate_technical_score_from_dict(audit_data)
    
    cursor.execute('''
        INSERT INTO seo_audits (url, audit_data, seo_score, content_score, technical_score, user_id)
        VALUES (?, ?, ?, ?, ?, ?)
    ''', (url, json.dumps(audit_data), seo_score, content_score, technical_score, user_id))
    
    # Also save to trends table
    cursor.execute('''
        INSERT INTO seo_trends (url, date, seo_score, content_score, technical_score, keywords_count, word_count, internal_links, external_links)
        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
    ''', (
        url, 
        datetime.now().date(), 
        seo_score, 
        content_score, 
        technical_score,
        len(audit_data.get('keywords', [])),
        audit_data.get('word_count', 0),
        audit_data.get('internal_links', 0),
        audit_data.get('external_links', 0)
    ))
    
    conn.commit()
    conn.close()

# Helper: Calculate SEO score based on metrics
def calculate_seo_score(audit_result: SEOAuditResponse) -> int:
    score = 0
    max_score = 100
    
    # Basic SEO factors (40 points)
    if audit_result.title_tag and len(audit_result.title_tag) <= 60:
        score += 10
    if audit_result.meta_description and len(audit_result.meta_description) <= 160:
        score += 10
    if audit_result.h1:
        score += 10
    if audit_result.keywords and len(audit_result.keywords) >= 5:
        score += 10
    
    # Technical factors (30 points)
    if audit_result.mobile_friendly:
        score += 10
    if audit_result.structured_data:
        score += 10
    if audit_result.canonical_tag:
        score += 5
    if audit_result.robots_meta:
        score += 5
    
    # Content factors (30 points)
    if audit_result.word_count and audit_result.word_count >= 300:
        score += 15
    if audit_result.internal_links and audit_result.internal_links >= 3:
        score += 10
    if audit_result.image_count and audit_result.image_count >= 2:
        score += 5
    
    return min(score, max_score)

def calculate_seo_score_from_dict(audit_data: dict) -> int:
    score = 0
    max_score = 100
    
    # Basic SEO factors (40 points)
    if audit_data.get('title_tag') and len(audit_data['title_tag']) <= 60:
        score += 10
    if audit_data.get('meta_description') and len(audit_data['meta_description']) <= 160:
        score += 10
    if audit_data.get('h1'):
        score += 10
    if audit_data.get('keywords') and len(audit_data['keywords']) >= 5:
        score += 10
    
    # Technical factors (30 points)
    if audit_data.get('mobile_friendly'):
        score += 10
    if audit_data.get('structured_data'):
        score += 10
    if audit_data.get('canonical_tag'):
        score += 5
    if audit_data.get('robots_meta'):
        score += 5
    
    # Content factors (30 points)
    if audit_data.get('word_count') and audit_data['word_count'] >= 300:
        score += 15
    if audit_data.get('internal_links') and audit_data['internal_links'] >= 3:
        score += 10
    if audit_data.get('image_count') and audit_data['image_count'] >= 2:
        score += 5
    
    return min(score, max_score)

# Helper: Calculate content quality score
def calculate_content_quality_score(audit_result: SEOAuditResponse) -> int:
    score = 0
    max_score = 100
    
    if audit_result.word_count:
        if audit_result.word_count >= 1000:
            score += 30
        elif audit_result.word_count >= 500:
            score += 20
        elif audit_result.word_count >= 300:
            score += 10
    
    if audit_result.readability:
        if "good" in audit_result.readability.lower() or "excellent" in audit_result.readability.lower():
            score += 25
        elif "fair" in audit_result.readability.lower():
            score += 15
    
    if audit_result.keywords and len(audit_result.keywords) >= 5:
        score += 20
    
    if audit_result.suggestions and len(audit_result.suggestions) <= 3:
        score += 15
    
    if not audit_result.missing_alt_text or "none" in audit_result.missing_alt_text.lower():
        score += 10
    
    return min(score, max_score)

def calculate_content_quality_score_from_dict(audit_data: dict) -> int:
    score = 0
    max_score = 100
    
    if audit_data.get('word_count'):
        if audit_data['word_count'] >= 1000:
            score += 30
        elif audit_data['word_count'] >= 500:
            score += 20
        elif audit_data['word_count'] >= 300:
            score += 10
    
    if audit_data.get('readability'):
        if "good" in audit_data['readability'].lower() or "excellent" in audit_data['readability'].lower():
            score += 25
        elif "fair" in audit_data['readability'].lower():
            score += 15
    
    if audit_data.get('keywords') and len(audit_data['keywords']) >= 5:
        score += 20
    
    if audit_data.get('suggestions') and len(audit_data['suggestions']) <= 3:
        score += 15
    
    if not audit_data.get('missing_alt_text') or "none" in audit_data.get('missing_alt_text', '').lower():
        score += 10
    
    return min(score, max_score)

# Helper: Calculate technical score
def calculate_technical_score(audit_result: SEOAuditResponse) -> int:
    score = 0
    max_score = 100
    
    if audit_result.mobile_friendly:
        score += 20
    if audit_result.structured_data:
        score += 20
    if audit_result.canonical_tag:
        score += 15
    if audit_result.robots_meta:
        score += 10
    if audit_result.social_meta:
        score += 15
    if audit_result.page_speed_score and audit_result.page_speed_score >= 80:
        score += 20
    
    return min(score, max_score)

def calculate_technical_score_from_dict(audit_data: dict) -> int:
    score = 0
    max_score = 100
    
    if audit_data.get('mobile_friendly'):
        score += 20
    if audit_data.get('structured_data'):
        score += 20
    if audit_data.get('canonical_tag'):
        score += 15
    if audit_data.get('robots_meta'):
        score += 10
    if audit_data.get('social_meta'):
        score += 15
    if audit_data.get('page_speed_score') and audit_data['page_speed_score'] >= 80:
        score += 20
    
    return min(score, max_score)

def get_url_hash(url):
    return hashlib.sha256(url.encode()).hexdigest()

async def async_crawl_website(start_url, max_pages=50, concurrency=10):
    visited = set()
    to_visit = [start_url]
    domain = urlparse(start_url).netloc
    urls = []
    sem = asyncio.Semaphore(concurrency)
    async with httpx.AsyncClient(timeout=10) as client:
        async def fetch(url):
            async with sem:
                try:
                    resp = await client.get(url)
                    resp.raise_for_status()
                    soup = BeautifulSoup(resp.text, 'html.parser')
                    links = []
                    for link in soup.find_all('a', href=True):
                        if not isinstance(link, Tag):
                            continue
                        href = link.get('href')
                        if not href or not isinstance(href, str):
                            continue
                        abs_url = urljoin(url, href)
                        parsed = urlparse(abs_url)
                        if parsed.netloc == domain and abs_url not in visited and abs_url not in to_visit:
                            if parsed.scheme in ['http', 'https']:
                                links.append(abs_url)
                    return url, links
                except Exception:
                    return url, []
        while to_visit and len(visited) < max_pages:
            batch = to_visit[:concurrency]
            to_visit = to_visit[concurrency:]
            results = await asyncio.gather(*(fetch(u) for u in batch))
            for url, found_links in results:
                if url not in visited:
                    visited.add(url)
                    urls.append(url)
                    for l in found_links:
                        if l not in visited and l not in to_visit and len(visited) + len(to_visit) < max_pages:
                            to_visit.append(l)
    return urls

async def async_check_broken_resources(soup, base_url, client):
    broken_images = []
    broken_links = []
    # Check images
    img_tasks = []
    for img in soup.find_all('img'):
        if not isinstance(img, Tag):
            continue
        src = img.get('src')
        if not src or not isinstance(src, str):
            continue
        # Convert relative URLs to absolute
        if src.startswith('//'):
            src = 'https:' + src
        elif src.startswith('/'):
            src = urljoin(base_url, src)
        elif not src.startswith(('http://', 'https://')):
            src = urljoin(base_url, src)
        img_tasks.append((src, client.head(src, timeout=5, follow_redirects=True)))
    # Check links
    link_tasks = []
    for link in soup.find_all('a', href=True):
        if not isinstance(link, Tag):
            continue
        href = link.get('href')
        if not href or not isinstance(href, str):
            continue
        if href.startswith(('mailto:', 'tel:', 'javascript:', '#')):
            continue
        if href.startswith('//'):
            href = 'https:' + href
        elif href.startswith('/'):
            href = urljoin(base_url, href)
        elif not href.startswith(('http://', 'https://')):
            href = urljoin(base_url, href)
        link_tasks.append((href, client.head(href, timeout=5, follow_redirects=True)))
    # Await all tasks
    img_results = await asyncio.gather(*[t[1] for t in img_tasks], return_exceptions=True)
    for (src, _), resp in zip(img_tasks, img_results):
        if isinstance(resp, Exception) or (hasattr(resp, 'status_code') and resp.status_code >= 400):
            broken_images.append(src)
    link_results = await asyncio.gather(*[t[1] for t in link_tasks], return_exceptions=True)
    for (href, _), resp in zip(link_tasks, link_results):
        if isinstance(resp, Exception) or (hasattr(resp, 'status_code') and resp.status_code >= 400):
            broken_links.append(href)
    return broken_images, broken_links

def check_broken_resources(soup, base_url):
    """Sync wrapper for async_check_broken_resources for legacy code paths."""
    import httpx
    async def runner():
        async with httpx.AsyncClient() as client:
            return await async_check_broken_resources(soup, base_url, client)
    return asyncio.run(runner())

async def async_extract_images_missing_alt(soup, base_url):
    images_missing_alt = []
    for img in soup.find_all('img'):
        if not isinstance(img, Tag):
            continue
        src = img.get('src')
        alt = img.get('alt', '')
        title = img.get('title', '')
        if not src or not isinstance(src, str):
            continue
        if src.startswith('//'):
            src = 'https:' + src
        elif src.startswith('/'):
            src = urljoin(base_url, src)
        elif not src.startswith(('http://', 'https://')):
            src = urljoin(base_url, src)
        if not alt or (isinstance(alt, str) and alt.strip() == ''):
            images_missing_alt.append({
                'src': src,
                'alt': alt or '',
                'title': title or '',
                'width': img.get('width', ''),
                'height': img.get('height', ''),
                'class': '',
                'id': img.get('id', '')
            })
    return images_missing_alt

def extract_images_missing_alt(soup, base_url):
    """Sync wrapper for async_extract_images_missing_alt for legacy code paths."""
    return asyncio.run(async_extract_images_missing_alt(soup, base_url))

# Helper: Extract all internal links
def extract_all_internal_links(soup, base_url):
    all_internal_links = []
    domain = urlparse(base_url).netloc
    
    for link in soup.find_all('a', href=True):
        if not isinstance(link, Tag):
            continue
        
        href = link.get('href')
        text = link.get_text(strip=True)
        title = link.get('title', '')
        
        if not href or not isinstance(href, str):
            continue
        
        # Skip mailto, tel, javascript, etc.
        if href.startswith(('mailto:', 'tel:', 'javascript:', '#')):
            continue
        
        # Convert relative URLs to absolute
        if href.startswith('//'):
            href = 'https:' + href
        elif href.startswith('/'):
            href = urljoin(base_url, href)
        elif not href.startswith(('http://', 'https://')):
            href = urljoin(base_url, href)
        
        # Check if it's an internal link
        parsed = urlparse(href)
        if parsed.netloc == domain:
            all_internal_links.append({
                'url': href,
                'text': text,
                'title': title,
                'class': '',
                'id': link.get('id', '')
            })
    
    return all_internal_links

# Helper: Extract all images
def extract_all_images(soup, base_url):
    all_images = []
    
    for img in soup.find_all('img'):
        if not isinstance(img, Tag):
            continue
        
        src = img.get('src')
        alt = img.get('alt', '')
        title = img.get('title', '')
        
        if not src or not isinstance(src, str):
            continue
        
        # Convert relative URLs to absolute
        if src.startswith('//'):
            src = 'https:' + src
        elif src.startswith('/'):
            src = urljoin(base_url, src)
        elif not src.startswith(('http://', 'https://')):
            src = urljoin(base_url, src)
        
        all_images.append({
            'src': src,
            'alt': alt or '',
            'title': title or '',
            'width': img.get('width', ''),
            'height': img.get('height', ''),
            'class': '',
            'id': img.get('id', ''),
            'has_alt': bool(alt and (isinstance(alt, str) and alt.strip()))
        })
    
    return all_images

# Helper: Generate sitemap XML
def generate_sitemap(urls):
    urlset = ET.Element('urlset', xmlns="http://www.sitemaps.org/schemas/sitemap/0.9")
    for url in urls:
        url_elem = ET.SubElement(urlset, 'url')
        loc = ET.SubElement(url_elem, 'loc')
        loc.text = url
    return ET.tostring(urlset, encoding='utf-8', method='xml')

# Background competitor analysis job
def run_competitor_analysis(job_id, target_url, competitor_urls):
    competitor_jobs[job_id]['status'] = 'analyzing'
    
    # Analyze target URL
    try:
        target_req = SEOAuditRequest(url=target_url)
        target_result = seo_audit(target_req)
        target_result.seo_score = calculate_seo_score(target_result)
        target_result.content_quality_score = calculate_content_quality_score(target_result)
        target_result.technical_score = calculate_technical_score(target_result)
        competitor_jobs[job_id]['target_analysis'] = target_result.dict()
    except Exception as e:
        competitor_jobs[job_id]['target_analysis'] = {'error': str(e)}
    
    # Analyze competitor URLs
    competitor_analyses = []
    for i, comp_url in enumerate(competitor_urls):
        try:
            comp_req = SEOAuditRequest(url=comp_url)
            comp_result = seo_audit(comp_req)
            comp_result.seo_score = calculate_seo_score(comp_result)
            comp_result.content_quality_score = calculate_content_quality_score(comp_result)
            comp_result.technical_score = calculate_technical_score(comp_result)
            competitor_analyses.append(comp_result.dict())
        except Exception as e:
            competitor_analyses.append({'url': comp_url, 'error': str(e)})
        
        competitor_jobs[job_id]['progress'] = f"Target: Done, Competitors: {i+1}/{len(competitor_urls)}"
    
    competitor_jobs[job_id]['competitor_analyses'] = competitor_analyses
    
    # Generate comparison metrics and recommendations
    if 'target_analysis' in competitor_jobs[job_id] and not 'error' in competitor_jobs[job_id]['target_analysis']:
        target = competitor_jobs[job_id]['target_analysis']
        comparison_metrics = {
            'target_seo_score': target.get('seo_score', 0),
            'target_content_score': target.get('content_quality_score', 0),
            'target_technical_score': target.get('technical_score', 0),
            'competitor_avg_seo': 0,
            'competitor_avg_content': 0,
            'competitor_avg_technical': 0,
            'rank_position': 1
        }
        
        valid_competitors = [c for c in competitor_analyses if 'error' not in c]
        if valid_competitors:
            avg_seo = sum(c.get('seo_score', 0) for c in valid_competitors) / len(valid_competitors)
            avg_content = sum(c.get('content_quality_score', 0) for c in valid_competitors) / len(valid_competitors)
            avg_technical = sum(c.get('technical_score', 0) for c in valid_competitors) / len(valid_competitors)
            
            comparison_metrics['competitor_avg_seo'] = round(avg_seo, 1)
            comparison_metrics['competitor_avg_content'] = round(avg_content, 1)
            comparison_metrics['competitor_avg_technical'] = round(avg_technical, 1)
            
            # Calculate rank position
            all_scores = [target.get('seo_score', 0)] + [c.get('seo_score', 0) for c in valid_competitors]
            all_scores.sort(reverse=True)
            comparison_metrics['rank_position'] = all_scores.index(target.get('seo_score', 0)) + 1
        
        competitor_jobs[job_id]['comparison_metrics'] = comparison_metrics
        
        # Generate AI-powered recommendations
        recommendations = generate_competitor_recommendations(target, valid_competitors)
        competitor_jobs[job_id]['recommendations'] = recommendations
    
    competitor_jobs[job_id]['status'] = 'done'

def generate_competitor_recommendations(target, competitors):
    if not competitors:
        return ["No competitor data available for comparison."]
    
    # Find the best competitor
    best_competitor = max(competitors, key=lambda x: x.get('seo_score', 0))
    
    recommendations = []
    
    # SEO Score recommendations
    if target.get('seo_score', 0) < best_competitor.get('seo_score', 0):
        recommendations.append(f"Improve overall SEO score: Your score is {target.get('seo_score', 0)} vs competitor's {best_competitor.get('seo_score', 0)}")
    
    # Content recommendations
    if target.get('content_quality_score', 0) < best_competitor.get('content_quality_score', 0):
        recommendations.append("Enhance content quality: Add more relevant keywords and improve readability")
    
    # Technical recommendations
    if target.get('technical_score', 0) < best_competitor.get('technical_score', 0):
        recommendations.append("Improve technical SEO: Add structured data, optimize mobile experience")
    
    # Specific recommendations based on missing elements
    if not target.get('mobile_friendly') and best_competitor.get('mobile_friendly'):
        recommendations.append("Add viewport meta tag for mobile optimization")
    
    if not target.get('structured_data') and best_competitor.get('structured_data'):
        recommendations.append("Implement structured data (JSON-LD) for better search visibility")
    
    if target.get('word_count', 0) and target.get('word_count', 0) < 500:
        recommendations.append("Increase content length: Aim for at least 500 words")
    
    if not target.get('social_meta') and best_competitor.get('social_meta'):
        recommendations.append("Add Open Graph and Twitter Card meta tags for social sharing")
    
    return recommendations[:5]  # Limit to 5 recommendations

# Background scan job
def run_full_scan(job_id, start_url):
    scan_jobs[job_id]['status'] = 'crawling'
    urls = async_crawl_website_sync(start_url)
    scan_jobs[job_id]['urls'] = urls
    scan_jobs[job_id]['status'] = 'auditing'
    results = []
    for idx, url in enumerate(urls):
        try:
            # Reuse the existing audit logic
            req = SEOAuditRequest(url=url)
            res = seo_audit(req)
            results.append({'url': url, 'result': res.dict()})
        except Exception as e:
            results.append({'url': url, 'error': str(e)})
        scan_jobs[job_id]['progress'] = f"{idx+1}/{len(urls)}"
    scan_jobs[job_id]['results'] = results
    scan_jobs[job_id]['status'] = 'done'
    scan_jobs[job_id]['sitemap'] = generate_sitemap(urls)

# Async wrapper for full scan
async def async_run_full_scan(job_id, start_url):
    scan_jobs[job_id]['status'] = 'crawling'
    urls = await async_crawl_website(start_url)
    scan_jobs[job_id]['urls'] = urls
    scan_jobs[job_id]['status'] = 'auditing'
    results = []
    for idx, url in enumerate(urls):
        try:
            # Reuse the existing audit logic
            req = SEOAuditRequest(url=url)
            res = seo_audit(req)
            results.append({'url': url, 'result': res.dict()})
        except Exception as e:
            results.append({'url': url, 'error': str(e)})
        scan_jobs[job_id]['progress'] = f"{idx+1}/{len(urls)}"
    scan_jobs[job_id]['results'] = results
    scan_jobs[job_id]['status'] = 'done'
    scan_jobs[job_id]['sitemap'] = generate_sitemap(urls)

# Synchronous wrapper for async crawl
def async_crawl_website_sync(start_url, max_pages=50):
    return asyncio.run(async_crawl_website(start_url, max_pages))

@app.post("/seo-audit", response_model=SEOAuditResponse)
def seo_audit(request: SEOAuditRequest):
    try:
        response = requests.get(request.url, timeout=10)
        response.raise_for_status()
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to fetch URL: {e}")

    soup = BeautifulSoup(response.text, 'html.parser')
    # Extract visible text
    texts = soup.stripped_strings
    content = ' '.join(texts)

    # --- Advanced metrics extraction ---
    # Internal/External links
    domain = urlparse(request.url).netloc
    internal_links = 0
    external_links = 0
    for a in soup.find_all('a', href=True):
        if not isinstance(a, Tag):
            continue
        href = a.get('href')
        if not href or not isinstance(href, str):
            continue
        abs_url = urljoin(request.url, href)
        parsed = urlparse(abs_url)
        if parsed.netloc == domain:
            internal_links += 1
        elif parsed.scheme in ['http', 'https']:
            external_links += 1
    # Image count
    image_count = len(soup.find_all('img'))
    # Mobile friendly (simple check for viewport meta)
    mobile_friendly = bool(soup.find('meta', attrs={'name': 'viewport'}))
    # Page speed (placeholder: use 80, or integrate with real API later)
    page_speed_score = 80
    # Structured data (look for JSON-LD or microdata)
    structured_data = bool(soup.find('script', type='application/ld+json')) or bool(soup.find(attrs={'itemscope': True}))
    # Canonical tag
    canonical_tag = bool(soup.find('link', rel='canonical'))
    # Robots meta
    robots_meta = bool(soup.find('meta', attrs={'name': 'robots'}))
    # Social meta (Open Graph or Twitter Card)
    social_meta = bool(soup.find('meta', attrs={'property': lambda v: isinstance(v, str) and v.startswith('og:')})) or bool(soup.find('meta', attrs={'name': 'twitter:card'}))
    
    # Check for broken images and links
    broken_images, broken_links = check_broken_resources(soup, request.url)
    
    # Extract images with missing alt text
    images_missing_alt = extract_images_missing_alt(soup, request.url)
    
    # Extract all internal links and images
    all_internal_links = extract_all_internal_links(soup, request.url)
    all_images = extract_all_images(soup, request.url)
    # --- End advanced metrics extraction ---

    # Enhanced prompt for DeepSeek
    prompt = f"""
You are an SEO expert. Analyze the following website content and provide:
1. 5-10 most relevant keywords (comma-separated).
2. A compelling meta description (max 160 characters).
3. An SEO-optimized title tag (max 60 characters).
4. The main H1 heading(s).
5. The total word count.
6. A readability score (Flesch-Kincaid or similar).
7. List any images missing alt text (if possible).
8. 3-5 actionable suggestions to improve SEO.
9. The single best focus keyphrase for this page.
10. 3-5 synonyms or close variants for the focus keyphrase (comma-separated).

Content:
{content[:3000]}

Format (respond exactly as below):
Keywords: ...
Meta Description: ...
Title Tag: ...
H1: ...
Word Count: ...
Readability: ...
Missing Alt Text: ...
Suggestions:
- ...
- ...
- ...
Focus Keyphrase: ...
Keyphrase Synonyms: ...
"""
    data = {
        "model": "deepseek-chat",
        "messages": [
            {"role": "system", "content": "You are an SEO expert."},
            {"role": "user", "content": prompt}
        ],
        "max_tokens": 600,
        "temperature": 0.7
    }
    headers = {
        "Authorization": f"Bearer {DEEPSEEK_API_KEY}",
        "Content-Type": "application/json"
    }
    print("[DEBUG] Authorization header:", headers["Authorization"])
    try:
        api_response = requests.post(DEEPSEEK_API_URL, headers=headers, json=data, timeout=30)
        print("[DEBUG] DeepSeek status code:", api_response.status_code)
        print("[DEBUG] DeepSeek response text:", api_response.text)
        api_response.raise_for_status()
        response_json = api_response.json()
        print("[DEBUG] DeepSeek full API response:", response_json)
        result = response_json["choices"][0]["message"]["content"] if "choices" in response_json and response_json["choices"] else ""
        print("[DEBUG] DeepSeek raw response:", result)

        # Regex extraction for all fields (more robust)
        def extract(pattern, text, default=None):
            m = re.search(pattern, text, re.IGNORECASE | re.DOTALL)
            return m.group(1).strip() if m else default

        keywords_str = extract(r"Keywords\s*:?\s*(.+)", result)
        meta_description = extract(r"Meta Description\s*:?\s*(.+)", result)
        title_tag = extract(r"Title Tag\s*:?\s*(.+)", result)
        h1 = extract(r"H1\s*:?\s*(.+)", result)
        word_count_str = extract(r"Word Count\s*:?\s*(\d+)", result)
        readability = extract(r"Readability\s*:?\s*(.+)", result)
        missing_alt_text = extract(r"Missing Alt Text\s*:?\s*(.+)", result)
        suggestions_block = extract(r"Suggestions\s*:?\s*([\s\S]+?)(?:Focus Keyphrase|Keyphrase Synonyms|$)", result)
        focus_keyphrase = extract(r"Focus Keyphrase\s*:?\s*(.+)", result, "")
        keyphrase_synonyms_str = extract(r"Keyphrase Synonyms\s*:?\s*(.+)", result, "")

        keywords = [k.strip() for k in keywords_str.split(",") if k.strip()] if keywords_str else []
        word_count = int(word_count_str) if word_count_str and word_count_str.isdigit() else None
        suggestions = []
        if suggestions_block:
            suggestions = [s.strip("- ").strip() for s in suggestions_block.split("\n") if s.strip() and s.strip().startswith("-")]
        keyphrase_synonyms = [k.strip() for k in keyphrase_synonyms_str.split(",") if k.strip()] if keyphrase_synonyms_str else []

        print("[DEBUG] extracted keywords:", keywords)
        print("[DEBUG] extracted meta_description:", meta_description)
        print("[DEBUG] extracted title_tag:", title_tag)
        print("[DEBUG] extracted h1:", h1)
        print("[DEBUG] extracted word_count:", word_count)
        print("[DEBUG] extracted readability:", readability)
        print("[DEBUG] extracted missing_alt_text:", missing_alt_text)
        print("[DEBUG] extracted suggestions:", suggestions)
        print("[DEBUG] extracted focus_keyphrase:", focus_keyphrase)
        print("[DEBUG] extracted keyphrase_synonyms:", keyphrase_synonyms)

        if not keywords or not meta_description:
            raise ValueError("Could not parse DeepSeek response.")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"DeepSeek API error: {e}")

    # After extracting focus_keyphrase and keyphrase_synonyms:
    from collections import Counter
    import string
    try:
        import nltk
        from nltk.corpus import stopwords
        from nltk.util import ngrams
        nltk.download('punkt', quiet=True)
        nltk.download('stopwords', quiet=True)
        stop_words = set(stopwords.words('english'))
        def clean_text(text):
            return text.lower().translate(str.maketrans('', '', string.punctuation))
        tokens = nltk.word_tokenize(content)
        tokens_clean = [t for t in tokens if t.isalnum() and t.lower() not in stop_words]
    except Exception:
        tokens = content.split()
        tokens_clean = tokens
        stop_words = set()
    # Keyphrase density
    keyphrase_density = None
    if focus_keyphrase:
        keyphrase_count = content.lower().count(focus_keyphrase.lower())
        keyphrase_density = keyphrase_count / (len(tokens) or 1)
    # Keyphrase placement
    keyphrase_placement = {}
    if focus_keyphrase:
        keyphrase_placement = {
            'in_title': focus_keyphrase.lower() in (title_tag or '').lower(),
            'in_meta': focus_keyphrase.lower() in (meta_description or '').lower(),
            'in_h1': focus_keyphrase.lower() in (h1 or '').lower(),
            'in_first_paragraph': focus_keyphrase.lower() in ' '.join(list(soup.stripped_strings)[:50]).lower()
        }
    # Synonym usage
    synonym_usage = {}
    if keyphrase_synonyms:
        for syn in keyphrase_synonyms:
            synonym_usage[syn] = content.lower().count(syn.lower())
    # LSI keywords (top 5 most common nouns not in keywords or synonyms)
    lsi_keywords = []
    try:
        from nltk import pos_tag
        tagged = pos_tag(tokens_clean)
        nouns = [w for w, pos in tagged if pos.startswith('NN')]
        lsi_candidates = [w for w in nouns if w not in keywords and (not keyphrase_synonyms or w not in keyphrase_synonyms)]
        lsi_keywords = [w for w, _ in Counter(lsi_candidates).most_common(5)]
    except Exception:
        lsi_keywords = []

    result = SEOAuditResponse(
        keywords=keywords,
        meta_description=meta_description,
        title_tag=title_tag,
        h1=h1,
        word_count=word_count,
        readability=readability,
        missing_alt_text=missing_alt_text,
        suggestions=suggestions,
        internal_links=internal_links,
        external_links=external_links,
        image_count=image_count,
        mobile_friendly=mobile_friendly,
        page_speed_score=page_speed_score,
        structured_data=structured_data,
        canonical_tag=canonical_tag,
        robots_meta=robots_meta,
        social_meta=social_meta,
        broken_images=broken_images,
        broken_links=broken_links,
        images_missing_alt=images_missing_alt,
        all_internal_links=all_internal_links,
        all_images=all_images,
        focus_keyphrase=focus_keyphrase,
        keyphrase_synonyms=keyphrase_synonyms,
        keyphrase_density=keyphrase_density,
        keyphrase_placement=keyphrase_placement,
        synonym_usage=synonym_usage,
        lsi_keywords=lsi_keywords
    )
    
    # Calculate scores and add them to the response
    result.seo_score = calculate_seo_score(result)
    result.content_quality_score = calculate_content_quality_score(result)
    result.technical_score = calculate_technical_score(result)
    
    # Save to database for analytics
    save_audit_to_db(request.url, result.dict())
    
    return result

@app.post("/competitor-analysis")
def competitor_analysis(request: CompetitorAnalysisRequest, background_tasks: BackgroundTasks):
    job_id = str(int(time.time() * 1000))
    competitor_jobs[job_id] = {
        'status': 'pending', 
        'progress': '0/0', 
        'target_analysis': {}, 
        'competitor_analyses': [], 
        'comparison_metrics': {},
        'recommendations': []
    }
    background_tasks.add_task(run_competitor_analysis, job_id, request.target_url, request.competitor_urls)
    return {"job_id": job_id}

@app.get("/competitor-status/{job_id}")
def competitor_status(job_id: str):
    job = competitor_jobs.get(job_id)
    if not job:
        return {"status": "not_found"}
    return {
        "status": job['status'],
        "progress": job['progress'],
        "target_analysis": job.get('target_analysis', {}),
        "competitor_analyses": job.get('competitor_analyses', []),
        "comparison_metrics": job.get('comparison_metrics', {}),
        "recommendations": job.get('recommendations', [])
    }

@app.post("/analytics/trends")
def get_analytics_trends(request: AnalyticsRequest):
    conn = sqlite3.connect('boostify_analytics.db')
    cursor = conn.cursor()
    
    # Calculate date range
    end_date = datetime.now().date()
    if request.date_range == "7d":
        start_date = end_date - timedelta(days=7)
    elif request.date_range == "30d":
        start_date = end_date - timedelta(days=30)
    elif request.date_range == "90d":
        start_date = end_date - timedelta(days=90)
    elif request.date_range == "1y":
        start_date = end_date - timedelta(days=365)
    else:
        start_date = end_date - timedelta(days=30)
    
    # Get trends data
    cursor.execute('''
        SELECT date, seo_score, content_score, technical_score, keywords_count, word_count, internal_links, external_links
        FROM seo_trends 
        WHERE url = ? AND date >= ?
        ORDER BY date ASC
    ''', (request.url, start_date))
    
    trends_data = cursor.fetchall()
    
    # Get summary statistics
    cursor.execute('''
        SELECT 
            AVG(seo_score) as avg_seo,
            AVG(content_score) as avg_content,
            AVG(technical_score) as avg_technical,
            MAX(seo_score) as max_seo,
            MIN(seo_score) as min_seo,
            COUNT(*) as total_audits
        FROM seo_trends 
        WHERE url = ? AND date >= ?
    ''', (request.url, start_date))
    
    summary = cursor.fetchone()
    
    conn.close()
    
    # Format trends data
    trends = []
    for row in trends_data:
        trends.append({
            "date": row[0],
            "seo_score": row[1],
            "content_score": row[2],
            "technical_score": row[3],
            "keywords_count": row[4],
            "word_count": row[5],
            "internal_links": row[6],
            "external_links": row[7]
        })
    
    return {
        "trends": trends,
        "summary": {
            "avg_seo": round(summary[0], 1) if summary[0] else 0,
            "avg_content": round(summary[1], 1) if summary[1] else 0,
            "avg_technical": round(summary[2], 1) if summary[2] else 0,
            "max_seo": summary[3] or 0,
            "min_seo": summary[4] or 0,
            "total_audits": summary[5] or 0
        }
    }

@app.get("/analytics/history/{url:path}")
def get_audit_history(url: str, limit: int = 10, offset: int = 0):
    conn = sqlite3.connect('boostify_analytics.db')
    cursor = conn.cursor()
    cursor.execute('''
        SELECT created_at, audit_data, seo_score, content_score, technical_score
        FROM seo_audits 
        WHERE url = ?
        ORDER BY created_at DESC
        LIMIT ? OFFSET ?
    ''', (url, limit, offset))
    history = cursor.fetchall()
    conn.close()
    return {
        "history": [
            {
                "created_at": row[0],
                "audit_data": json.loads(row[1]),
                "seo_score": row[2],
                "content_score": row[3],
                "technical_score": row[4]
            }
            for row in history
        ]
    }

@app.get("/analytics/urls")
def get_all_scanned_urls():
    conn = sqlite3.connect('boostify_analytics.db')
    cursor = conn.cursor()
    cursor.execute('''SELECT DISTINCT url FROM seo_audits ORDER BY url ASC''')
    urls = [row[0] for row in cursor.fetchall()]
    conn.close()
    return {"urls": urls}

@app.get("/analytics/export/{url:path}")
def export_audit_history_csv(url: str, format: str = "csv", limit: int = 100):
    conn = sqlite3.connect('boostify_analytics.db')
    cursor = conn.cursor()
    cursor.execute('''
        SELECT created_at, audit_data, seo_score, content_score, technical_score
        FROM seo_audits 
        WHERE url = ?
        ORDER BY created_at DESC
        LIMIT ?
    ''', (url, limit))
    rows = cursor.fetchall()
    conn.close()
    output = io.StringIO()
    writer = csv.writer(output)
    writer.writerow(["created_at", "seo_score", "content_score", "technical_score", "audit_data"])
    for row in rows:
        writer.writerow([row[0], row[2], row[3], row[4], row[1]])
    output.seek(0)
    if format == "csv":
        return StreamingResponse(iter([output.getvalue()]), media_type="text/csv", headers={"Content-Disposition": f"attachment; filename=seo_audit_{url.replace('/', '_')}.csv"})
    elif format == "excel":
        import pandas as pd
        df = pd.read_csv(io.StringIO(output.getvalue()))
        excel_output = io.BytesIO()
        df.to_excel(excel_output, index=False)
        excel_output.seek(0)
        return StreamingResponse(excel_output, media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", headers={"Content-Disposition": f"attachment; filename=seo_audit_{url.replace('/', '_')}.xlsx"})
    else:
        return {"error": "Invalid format"}

@app.post("/full-scan")
def full_scan(request: SEOAuditRequest, background_tasks: BackgroundTasks):
    job_id = str(int(time.time() * 1000))
    scan_jobs[job_id] = {'status': 'pending', 'progress': '0/0', 'urls': [], 'results': [], 'sitemap': b''}
    background_tasks.add_task(lambda: asyncio.run(async_run_full_scan(job_id, request.url)))
    return {"job_id": job_id}

@app.get("/scan-status/{job_id}")
def scan_status(job_id: str):
    job = scan_jobs.get(job_id)
    if not job:
        return {"status": "not_found"}
    return {
        "status": job['status'],
        "progress": job['progress'],
        "urls": job['urls'],
        "results": job['results'] if job['status'] == 'done' else []
    }

@app.get("/sitemap.xml/{job_id}")
def sitemap_xml(job_id: str):
    job = scan_jobs.get(job_id)
    if not job or not job.get('sitemap'):
        return StreamingResponse(iter([b'']), media_type="application/xml")
    return StreamingResponse(iter([job['sitemap']]), media_type="application/xml")

# Enhanced Export Endpoints
@app.post("/export/report")
def export_report(request: ExportRequest):
    try:
        # Get audit data
        audit_req = SEOAuditRequest(url=request.url)
        audit_result = seo_audit(audit_req)
        audit_data = audit_result.dict()
        
        # Generate report based on format
        template = request.template or "default"
        if request.format == "pdf":
            report_data = generate_pdf_report(audit_data, template)
            return StreamingResponse(
                iter([report_data]), 
                media_type="application/pdf",
                headers={"Content-Disposition": f"attachment; filename=seo-audit-{request.url.replace('://', '-').replace('/', '-')}.pdf"}
            )
        elif request.format == "pptx":
            report_data = generate_pptx_report(audit_data, template)
            return StreamingResponse(
                iter([report_data]), 
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                headers={"Content-Disposition": f"attachment; filename=seo-audit-{request.url.replace('://', '-').replace('/', '-')}.pptx"}
            )
        elif request.format == "html":
            report_data = generate_html_report(audit_data, template)
            return StreamingResponse(
                iter([report_data.encode('utf-8')]), 
                media_type="text/html",
                headers={"Content-Disposition": f"attachment; filename=seo-audit-{request.url.replace('://', '-').replace('/', '-')}.html"}
            )
        elif request.format == "json":
            return audit_result
        else:
            raise HTTPException(status_code=400, detail="Unsupported format")
            
    except Exception as e:
        import traceback
        tb = traceback.format_exc()
        print(f"[EXPORT ERROR] {e}\n{tb}")
        raise HTTPException(status_code=500, detail=f"Export failed: {str(e)}\n{tb}")

@app.post("/export/schedule")
def schedule_export_report(request: ScheduledReportRequest):
    try:
        report_id = f"report_{int(time.time() * 1000)}"
        
        # Schedule the report
        schedule_report(
            report_id=report_id,
            url=request.url,
            frequency=request.frequency,
            format=request.format,
            template=request.template,
            email=request.email
        )
        
        return {
            "report_id": report_id,
            "status": "scheduled",
            "frequency": request.frequency,
            "format": request.format,
            "message": f"Report scheduled for {request.frequency} delivery"
        }
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Scheduling failed: {str(e)}")

@app.get("/export/scheduled/{report_id}")
def get_scheduled_report(report_id: str):
    try:
        conn = sqlite3.connect('boostify_analytics.db')
        cursor = conn.cursor()
        
        cursor.execute('''
            SELECT report_data, format, created_at FROM scheduled_reports 
            WHERE report_id = ? ORDER BY created_at DESC LIMIT 1
        ''', (report_id,))
        
        result = cursor.fetchone()
        conn.close()
        
        if not result:
            raise HTTPException(status_code=404, detail="Report not found")
        
        report_data, format_type, created_at = result
        
        if format_type == "pdf":
            media_type = "application/pdf"
            filename = f"scheduled-report-{report_id}.pdf"
        elif format_type == "pptx":
            media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            filename = f"scheduled-report-{report_id}.pptx"
        elif format_type == "html":
            media_type = "text/html"
            filename = f"scheduled-report-{report_id}.html"
        else:
            media_type = "application/octet-stream"
            filename = f"scheduled-report-{report_id}"
        
        return StreamingResponse(
            iter([report_data]), 
            media_type=media_type,
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to retrieve report: {str(e)}")

@app.get("/export/templates")
def get_export_templates():
    return {
        "templates": [
            {
                "id": "default",
                "name": "Default Report",
                "description": "Standard SEO audit report with all metrics",
                "formats": ["pdf", "pptx", "html", "json"]
            },
            {
                "id": "executive",
                "name": "Executive Summary",
                "description": "High-level summary for executives",
                "formats": ["pdf", "pptx", "html"]
            },
            {
                "id": "detailed",
                "name": "Detailed Analysis",
                "description": "Comprehensive report with technical details",
                "formats": ["pdf", "html", "json"]
            }
        ]
    }

# AI Content Writer Models
class ContentGenerationRequest(BaseModel):
    content_type: str  # "blog-post", "product-description", "landing-page", "meta-content"
    keywords: str
    length: str  # "short", "medium", "long", "comprehensive"
    tone: str  # "professional", "casual", "friendly", "authoritative", "conversational"
    audience: str  # "general", "professionals", "beginners", "experts", "business"
    audit_results: Optional[dict] = None

class ContentGenerationResponse(BaseModel):
    content: str
    seo_optimized: bool
    word_count: int
    keywords_used: List[str]

class ContentOptimizationRequest(BaseModel):
    url: str
    content_type: str = "landing-page"
    focus_keywords: str = ""
    target_audience: str = ""

class MetaTagRequest(BaseModel):
    url: str
    keywords: str
    page_type: str = "homepage"

class ImplementationPlanRequest(BaseModel):
    url: str
    priority_level: str = "medium"  # low, medium, high
    focus_areas: List[str] = []

class HTMLOptimizationRequest(BaseModel):
    url: str
    optimization_focus: str = "content"  # content, structure, accessibility

@app.get("/ai-content-writer.html")
def read_ai_content_writer():
    return FileResponse("ai-content-writer.html")

@app.get("/seo-implementation.html")
def read_seo_implementation():
    return FileResponse("seo-implementation.html")

@app.post("/ai-content-writer/generate")
def generate_content(request: ContentGenerationRequest):
    try:
        # Build the prompt based on content type and parameters
        prompt = build_content_prompt(request)
        
        # Call DeepSeek API for content generation
        data = {
            "model": "deepseek-chat",
            "messages": [
                {"role": "system", "content": "You are an expert SEO content writer. Create high-quality, SEO-optimized content that ranks well in search engines."},
                {"role": "user", "content": prompt}
            ],
            "max_tokens": 2000,
            "temperature": 0.7
        }
        
        headers = {
            "Authorization": f"Bearer {DEEPSEEK_API_KEY}",
            "Content-Type": "application/json"
        }
        
        api_response = requests.post(DEEPSEEK_API_URL, headers=headers, json=data, timeout=60)
        api_response.raise_for_status()
        response_json = api_response.json()
        
        generated_content = response_json["choices"][0]["message"]["content"]
        
        # Format the generated content for better readability
        formatted_content = format_generated_content(generated_content, request.content_type)
        
        # Extract keywords used and count words
        keywords_used = extract_keywords_from_content(formatted_content, request.keywords)
        word_count = len(formatted_content.split())
        
        return ContentGenerationResponse(
            content=formatted_content,
            seo_optimized=True,
            word_count=word_count,
            keywords_used=keywords_used
        )
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Content generation failed: {str(e)}")

def build_content_prompt(request: ContentGenerationRequest) -> str:
    content_type_prompts = {
        "blog-post": "Write an engaging blog post that provides valuable information to readers with proper structure and formatting.",
        "product-description": "Write a compelling product description that highlights benefits and features with clear formatting.",
        "landing-page": "Write conversion-focused landing page content that drives action with proper sections and formatting.",
        "meta-content": "Write optimized meta titles and descriptions for search engines with proper formatting."
    }
    
    length_guidelines = {
        "short": "300-500 words",
        "medium": "500-1000 words", 
        "long": "1000-2000 words",
        "comprehensive": "2000+ words"
    }
    
    tone_guidelines = {
        "professional": "Use formal, business-appropriate language",
        "casual": "Use relaxed, conversational language",
        "friendly": "Use warm, approachable language",
        "authoritative": "Use confident, expert-level language",
        "conversational": "Use natural, dialogue-like language"
    }
    
    audience_guidelines = {
        "general": "Write for a general audience with varying knowledge levels",
        "professionals": "Write for industry professionals with technical knowledge",
        "beginners": "Write for newcomers to the topic, explaining basic concepts",
        "experts": "Write for subject matter experts with deep knowledge",
        "business": "Write for business decision makers and executives"
    }
    
    # Build SEO recommendations from audit results
    seo_recommendations = ""
    if request.audit_results:
        audit = request.audit_results
        seo_recommendations = f"""
SEO Recommendations from Audit:
- Keywords to target: {', '.join(audit.get('keywords', []))}
- Meta description: {audit.get('meta_description', 'N/A')}
- Title tag: {audit.get('title_tag', 'N/A')}
- Content suggestions: {', '.join(audit.get('suggestions', []))}
- SEO Score: {audit.get('seo_score', 'N/A')}
- Content Quality Score: {audit.get('content_quality_score', 'N/A')}
- Technical Score: {audit.get('technical_score', 'N/A')}
"""
    
    prompt = f"""
{content_type_prompts.get(request.content_type, "Write SEO-optimized content with proper formatting.")}

Content Requirements:
- Type: {request.content_type.replace('-', ' ').title()}
- Target Keywords: {request.keywords}
- Length: {length_guidelines.get(request.length, "500-1000 words")}
- Tone: {tone_guidelines.get(request.tone, "Professional")}
- Audience: {audience_guidelines.get(request.audience, "General")}

Formatting Requirements:
- Use clear, well-structured headings (H1, H2, H3)
- Include proper paragraph breaks for readability
- Use bullet points and numbered lists where appropriate
- Include a compelling introduction and conclusion
- Use subheadings to break up content into logical sections
- Ensure proper spacing and formatting for easy reading

SEO Requirements:
- Naturally incorporate target keywords throughout the content
- Use proper heading hierarchy (H1 for main title, H2 for sections, H3 for subsections)
- Include internal linking opportunities
- Optimize for featured snippets with clear, concise answers
- Ensure readability with short paragraphs and clear structure
- Include a meta description and title tag if applicable

{seo_recommendations}

Please create high-quality, SEO-optimized content that follows these guidelines. Make it engaging, informative, and valuable to the target audience while maintaining strong SEO practices and excellent readability.
"""
    
    return prompt

def extract_keywords_from_content(content: str, target_keywords: str) -> List[str]:
    """Extract which target keywords were used in the generated content"""
    target_keyword_list = [kw.strip().lower() for kw in target_keywords.split(',')]
    content_lower = content.lower()
    used_keywords = []
    
    for keyword in target_keyword_list:
        if keyword in content_lower:
            used_keywords.append(keyword)
    
    return used_keywords

def format_generated_content(content: str, content_type: str) -> str:
    """Format the generated content for better readability"""
    
    # Clean up the content
    content = content.strip()
    
    # Add proper spacing around headings
    content = re.sub(r'^(#+)\s*(.+)$', r'\1 \2', content, flags=re.MULTILINE)
    
    # Ensure proper spacing after headings
    content = re.sub(r'^(#+)\s*(.+)$', r'\1 \2\n', content, flags=re.MULTILINE)
    
    # Add spacing around bullet points and numbered lists
    content = re.sub(r'^(\s*[-*•]\s+)', r'\n\1', content, flags=re.MULTILINE)
    content = re.sub(r'^(\s*\d+\.\s+)', r'\n\1', content, flags=re.MULTILINE)
    
    # Ensure proper paragraph spacing
    content = re.sub(r'\n{3,}', r'\n\n', content)
    
    # Format based on content type
    if content_type == "blog-post":
        # Ensure blog posts have proper structure
        if not content.startswith('# '):
            # Add a main heading if none exists
            lines = content.split('\n')
            if lines and lines[0].strip():
                content = f"# {lines[0].strip()}\n\n" + '\n'.join(lines[1:])
    
    elif content_type == "product-description":
        # Format product descriptions with clear sections
        sections = [
            "## Product Overview",
            "## Key Features",
            "## Benefits",
            "## Technical Specifications",
            "## Why Choose This Product"
        ]
        
        # Add sections if they don't exist
        for section in sections:
            if section not in content:
                content += f"\n\n{section}\n\n"
    
    elif content_type == "landing-page":
        # Format landing pages with conversion-focused sections
        sections = [
            "## Hero Section",
            "## Key Benefits",
            "## Features",
            "## Social Proof",
            "## Call to Action"
        ]
        
        # Add sections if they don't exist
        for section in sections:
            if section not in content:
                content += f"\n\n{section}\n\n"
    
    elif content_type == "meta-content":
        # Format meta content with clear structure
        if "Title:" not in content and "Meta Description:" not in content:
            content = f"## Meta Title\n\n[Your optimized title here]\n\n## Meta Description\n\n[Your optimized description here]\n\n## Additional Meta Tags\n\n[Any additional meta tags]"
    
    # Final cleanup
    content = re.sub(r'\n{3,}', r'\n\n', content)
    content = content.strip()
    
    return content

@app.delete("/export/schedule/{report_id}")
def cancel_scheduled_report(report_id: str):
    try:
        if report_id in scheduled_reports:
            # Remove from scheduled reports
            del scheduled_reports[report_id]
            
            # Remove from database
            conn = sqlite3.connect('boostify_analytics.db')
            cursor = conn.cursor()
            cursor.execute('DELETE FROM scheduled_reports WHERE report_id = ?', (report_id,))
            conn.commit()
            conn.close()
            
            return {"message": f"Report {report_id} cancelled successfully"}
        else:
            raise HTTPException(status_code=404, detail="Scheduled report not found")
            
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to cancel report: {str(e)}")

# Helper: Generate PDF report
def generate_pdf_report(audit_data: dict, template: str = "default") -> bytes:
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4)
    styles = getSampleStyleSheet()
    story = []
    
    # Title
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        spaceAfter=30,
        alignment=1,  # Center
        textColor=colors.HexColor('#1e40af')
    )
    story.append(Paragraph("SEO Audit Report", title_style))
    story.append(Spacer(1, 20))
    
    # Executive Summary
    story.append(Paragraph("Executive Summary", styles['Heading2']))
    story.append(Spacer(1, 12))
    
    # Calculate scores
    seo_score = calculate_seo_score_from_dict(audit_data)
    content_score = calculate_content_quality_score_from_dict(audit_data)
    technical_score = calculate_technical_score_from_dict(audit_data)
    
    summary_data = [
        ['Metric', 'Score', 'Status'],
        ['SEO Score', f'{seo_score}/100', 'Excellent' if seo_score >= 80 else 'Good' if seo_score >= 60 else 'Needs Improvement'],
        ['Content Quality', f'{content_score}/100', 'Excellent' if content_score >= 80 else 'Good' if content_score >= 60 else 'Needs Improvement'],
        ['Technical Score', f'{technical_score}/100', 'Excellent' if technical_score >= 80 else 'Good' if technical_score >= 60 else 'Needs Improvement'],
    ]
    
    summary_table = Table(summary_data)
    summary_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#1e40af')),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 12),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
        ('GRID', (0, 0), (-1, -1), 1, colors.black)
    ]))
    story.append(summary_table)
    story.append(Spacer(1, 20))
    
    # Key Metrics
    story.append(Paragraph("Key Metrics", styles['Heading2']))
    story.append(Spacer(1, 12))
    
    metrics_data = [
        ['Metric', 'Value'],
        ['Keywords Found', str(len(audit_data.get('keywords', [])))],
        ['Word Count', str(audit_data.get('word_count', 0))],
        ['Meta Description', 'Present' if audit_data.get('meta_description') else 'Missing'],
        ['Title Tag', 'Present' if audit_data.get('title_tag') else 'Missing'],
        ['H1 Tag', 'Present' if audit_data.get('h1') else 'Missing'],
        ['Internal Links', str(audit_data.get('internal_links', 0))],
        ['External Links', str(audit_data.get('external_links', 0))],
        ['Images', str(audit_data.get('image_count', 0))],
    ]
    
    metrics_table = Table(metrics_data)
    metrics_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#3b82f6')),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 12),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -1), colors.white),
        ('GRID', (0, 0), (-1, -1), 1, colors.black)
    ]))
    story.append(metrics_table)
    story.append(Spacer(1, 20))
    
    # Keywords
    if audit_data.get('keywords'):
        story.append(Paragraph("Keywords", styles['Heading2']))
        story.append(Spacer(1, 12))
        keywords_text = ", ".join(audit_data['keywords'])
        story.append(Paragraph(keywords_text, styles['Normal']))
        story.append(Spacer(1, 20))

    # Focus Keyphrase
    if audit_data.get('focus_keyphrase'):
        story.append(Paragraph("Focus Keyphrase", styles['Heading2']))
        story.append(Spacer(1, 12))
        story.append(Paragraph(audit_data['focus_keyphrase'], styles['Normal']))
        story.append(Spacer(1, 12))
        if audit_data.get('keyphrase_synonyms'):
            story.append(Paragraph("Keyphrase Synonyms", styles['Heading3']))
            story.append(Spacer(1, 8))
            synonyms_text = ", ".join(audit_data['keyphrase_synonyms'])
            story.append(Paragraph(synonyms_text, styles['Normal']))
            story.append(Spacer(1, 20))

    # Recommendations
    if audit_data.get('suggestions'):
        story.append(Paragraph("Recommendations", styles['Heading2']))
        story.append(Spacer(1, 12))
        for i, suggestion in enumerate(audit_data['suggestions'], 1):
            story.append(Paragraph(f"{i}. {suggestion}", styles['Normal']))
            story.append(Spacer(1, 6))
    
    doc.build(story)
    buffer.seek(0)
    return buffer.getvalue()

# Helper: Generate PowerPoint presentation
def generate_pptx_report(audit_data: dict, template: str = "default") -> bytes:
    if not PPTX_AVAILABLE:
        raise ImportError("PPTX library is not installed. Please install it with 'pip install python-pptx'")
    
    prs = Presentation()
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    if title:
        title.text = "SEO Audit Report"
    if subtitle and hasattr(subtitle, 'text'):
        subtitle.text = f"Analysis of {audit_data.get('url', 'Website')}"
    
    # Score slide
    score_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(score_slide_layout)
    title = slide.shapes.title
    
    # Calculate scores
    seo_score = calculate_seo_score_from_dict(audit_data)
    content_score = calculate_content_quality_score_from_dict(audit_data)
    technical_score = calculate_technical_score_from_dict(audit_data)
    
    if title:
        title.text = "SEO Performance Scores"
    
    # Add score content
    content = slide.placeholders[1]
    if hasattr(content, 'text_frame'):
        tf = content.text_frame
        tf.text = f"SEO Score: {seo_score}/100\nContent Quality: {content_score}/100\nTechnical Score: {technical_score}/100"
    
    # Metrics slide
    metrics_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(metrics_slide_layout)
    title = slide.shapes.title
    if title:
        title.text = "Key Metrics"
    
    content = slide.placeholders[1]
    if hasattr(content, 'text_frame'):
        tf = content.text_frame
        tf.text = f"Keywords: {len(audit_data.get('keywords', []))}\nWord Count: {audit_data.get('word_count', 0)}\nInternal Links: {audit_data.get('internal_links', 0)}\nExternal Links: {audit_data.get('external_links', 0)}"
    
    # Recommendations slide
    if audit_data.get('suggestions'):
        rec_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(rec_slide_layout)
        title = slide.shapes.title
        if title:
            title.text = "Recommendations"
        
        content = slide.placeholders[1]
        if hasattr(content, 'text_frame'):
            tf = content.text_frame
            tf.text = "\n".join([f"• {suggestion}" for suggestion in audit_data['suggestions']])
    
    # Save to buffer
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()

# Helper: Generate HTML report
def generate_html_report(audit_data: dict, template: str = "default") -> str:
    seo_score = calculate_seo_score_from_dict(audit_data)
    content_score = calculate_content_quality_score_from_dict(audit_data)
    technical_score = calculate_technical_score_from_dict(audit_data)
    
    html_template = f"""
    <!DOCTYPE html>
    <html lang="en">
    <head>
        <meta charset="UTF-8">
        <meta name="viewport" content="width=device-width, initial-scale=1.0">
        <title>SEO Audit Report</title>
        <style>
            body {{ font-family: Arial, sans-serif; margin: 40px; }}
            .header {{ background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); color: white; padding: 30px; border-radius: 10px; text-align: center; }}
            .score-card {{ background: #f8f9fa; padding: 20px; border-radius: 10px; margin: 20px 0; }}
            .metric {{ margin: 10px 0; padding: 10px; background: white; border-radius: 5px; border-left: 4px solid #3b82f6; }}
            .recommendation {{ background: #e3f2fd; padding: 15px; margin: 10px 0; border-radius: 5px; }}
            table {{ width: 100%; border-collapse: collapse; margin: 20px 0; }}
            th, td {{ padding: 12px; text-align: left; border-bottom: 1px solid #ddd; }}
            th {{ background-color: #3b82f6; color: white; }}
        </style>
    </head>
    <body>
        <div class="header">
            <h1>SEO Audit Report</h1>
            <p>Generated on {datetime.now().strftime('%B %d, %Y')}</p>
        </div>
        
        <div class="score-card">
            <h2>Performance Scores</h2>
            <table>
                <tr><th>Metric</th><th>Score</th><th>Status</th></tr>
                <tr><td>SEO Score</td><td>{seo_score}/100</td><td>{'Excellent' if seo_score >= 80 else 'Good' if seo_score >= 60 else 'Needs Improvement'}</td></tr>
                <tr><td>Content Quality</td><td>{content_score}/100</td><td>{'Excellent' if content_score >= 80 else 'Good' if content_score >= 60 else 'Needs Improvement'}</td></tr>
                <tr><td>Technical Score</td><td>{technical_score}/100</td><td>{'Excellent' if technical_score >= 80 else 'Good' if technical_score >= 60 else 'Needs Improvement'}</td></tr>
            </table>
        </div>
        
        <div class="metric">
            <h3>Key Metrics</h3>
            <p><strong>Keywords:</strong> {len(audit_data.get('keywords', []))}</p>
            <p><strong>Word Count:</strong> {audit_data.get('word_count', 0)}</p>
            <p><strong>Meta Description:</strong> {'Present' if audit_data.get('meta_description') else 'Missing'}</p>
            <p><strong>Title Tag:</strong> {'Present' if audit_data.get('title_tag') else 'Missing'}</p>
            <p><strong>Internal Links:</strong> {audit_data.get('internal_links', 0)}</p>
            <p><strong>External Links:</strong> {audit_data.get('external_links', 0)}</p>
        </div>
        
        {f'<div class="recommendation"><h3>Recommendations</h3><ul>' + ''.join([f'<li>{suggestion}</li>' for suggestion in audit_data.get('suggestions', [])]) + '</ul></div>' if audit_data.get('suggestions') else ''}
    </body>
    </html>
    """
    
    return html_template

# Helper: Schedule automated reports
scheduled_reports = {}

def schedule_report(report_id: str, url: str, frequency: str, format: str, template: str, email: Optional[str] = None):
    if not SCHEDULE_AVAILABLE:
        raise ImportError("Schedule library is not installed. Please install it with 'pip install schedule'")

    def run_scheduled_report():
        try:
            # Run audit
            req = SEOAuditRequest(url=url)
            result = seo_audit(req)
            
            # Generate report
            if format == "pdf":
                report_data = generate_pdf_report(result.dict(), template)
            elif format == "pptx":
                report_data = generate_pptx_report(result.dict(), template)
            elif format == "html":
                report_data = generate_html_report(result.dict(), template)
            
            # Save report to database
            conn = sqlite3.connect('boostify_analytics.db')
            cursor = conn.cursor()
            cursor.execute('''
                INSERT INTO scheduled_reports (report_id, url, format, report_data, created_at)
                VALUES (?, ?, ?, ?, ?)
            ''', (report_id, url, format, report_data, datetime.now()))
            conn.commit()
            conn.close()
            
            print(f"Scheduled report {report_id} completed for {url}")
            
        except Exception as e:
            print(f"Error in scheduled report {report_id}: {e}")
    
    if frequency == "daily":
        schedule.every().day.at("09:00").do(run_scheduled_report)
    elif frequency == "weekly":
        schedule.every().monday.at("09:00").do(run_scheduled_report)
    elif frequency == "monthly":
        # Schedule for first day of each month at 9 AM
        schedule.every().day.at("09:00").do(run_scheduled_report).tag(f"monthly_{report_id}")
    
    scheduled_reports[report_id] = {
        'url': url,
        'frequency': frequency,
        'format': format,
        'template': template,
        'email': email
    } 

@app.post("/seo-implementation/generate-content")
def generate_optimized_content(request: ContentOptimizationRequest):
    """Generate optimized content for a website based on SEO audit results"""
    try:
        # First run an SEO audit to get current issues
        audit_result = run_seo_audit(request.url)
        
        # Generate optimized content based on audit findings
        optimized_content = generate_content_from_audit(audit_result, request.content_type)
        
        return {
            "original_audit": audit_result,
            "optimized_content": optimized_content,
            "implementation_guide": generate_implementation_guide(audit_result, optimized_content)
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Content generation failed: {str(e)}")

@app.post("/seo-implementation/generate-meta-tags")
async def generate_optimized_meta_tags(request: MetaTagRequest):
    """Generate optimized meta tags for a website"""
    try:
        # Analyze current meta tags
        current_meta = await analyze_current_meta_tags(request.url)
        
        # Generate optimized versions
        optimized_meta = await generate_optimized_meta_tags_content(current_meta, request.keywords)
        
        return {
            "current_meta": current_meta,
            "optimized_meta": optimized_meta,
            "implementation_code": generate_meta_implementation_code(optimized_meta)
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Meta tag generation failed: {str(e)}")

@app.post("/seo-implementation/create-implementation-plan")
async def create_seo_implementation_plan(request: ImplementationPlanRequest):
    """Create a comprehensive SEO implementation plan"""
    try:
        # Run comprehensive audit
        audit_result = await run_comprehensive_audit(request.url)
        
        # Generate implementation plan
        plan = await generate_implementation_plan(audit_result, request.priority_level)
        
        return {
            "audit_results": audit_result,
            "implementation_plan": plan,
            "timeline": generate_implementation_timeline(plan),
            "resources_needed": generate_resource_requirements(plan)
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Implementation plan generation failed: {str(e)}")

@app.post("/seo-implementation/generate-html-code")
async def generate_optimized_html_code(request: HTMLOptimizationRequest):
    """Generate optimized HTML code for a website"""
    try:
        # Analyze current HTML structure
        current_html = await analyze_html_structure(request.url)
        
        # Generate optimized HTML
        optimized_html = await generate_optimized_html(current_html, request.optimization_focus)
        
        return {
            "current_html_analysis": current_html,
            "optimized_html": optimized_html,
            "changes_summary": generate_changes_summary(current_html, optimized_html)
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"HTML optimization failed: {str(e)}")

# SEO Implementation Helper Functions
def run_seo_audit(url: str) -> dict:
    """Run SEO audit and return results"""
    try:
        request = SEOAuditRequest(url=url)
        result = seo_audit(request)
        return result.dict()
    except Exception as e:
        raise Exception(f"SEO audit failed: {str(e)}")

def generate_content_from_audit(audit_result: dict, content_type: str) -> dict:
    """Generate optimized content based on audit results"""
    try:
        # Extract key information from audit
        keywords = audit_result.get('keywords', [])
        suggestions = audit_result.get('suggestions', [])
        seo_score = audit_result.get('seo_score', 0)
        
        # Create content generation request
        content_request = ContentGenerationRequest(
            content_type=content_type,
            keywords=", ".join(keywords[:5]),  # Use top 5 keywords
            tone="professional",
            length="comprehensive",
            audience="business",
            audit_results=audit_result
        )
        
        # Generate content
        content_response = generate_content(content_request)
        
        return {
            "content": content_response.content,
            "word_count": content_response.word_count,
            "keywords_used": content_response.keywords_used,
            "seo_optimized": content_response.seo_optimized
        }
    except Exception as e:
        raise Exception(f"Content generation failed: {str(e)}")

def generate_implementation_guide(audit_result: dict, optimized_content: dict) -> dict:
    """Generate implementation guide for SEO improvements"""
    return {
        "priority_actions": [
            "Replace current homepage content with optimized version",
            "Update meta tags with generated versions",
            "Fix broken links identified in audit",
            "Add missing alt text to images",
            "Improve heading structure"
        ],
        "content_improvements": [
            "Increase word count from current ~100 to 500+ words",
            "Add proper H1, H2, H3 heading hierarchy",
            "Include target keywords naturally throughout content",
            "Add internal linking opportunities",
            "Improve readability with shorter paragraphs"
        ],
        "technical_improvements": [
            "Fix broken links and images",
            "Add missing alt text",
            "Optimize meta descriptions",
            "Improve page speed",
            "Add structured data markup"
        ],
        "timeline": {
            "week_1": "Content structure and meta tag updates",
            "week_2": "Technical fixes and broken link resolution",
            "week_3": "Content expansion and optimization",
            "week_4": "Monitoring and fine-tuning"
        }
    }

async def analyze_current_meta_tags(url: str) -> dict:
    """Analyze current meta tags on a website"""
    try:
        response = requests.get(url, timeout=10)
        soup = BeautifulSoup(response.content, 'html.parser')
        
        meta_tags = {
            "title": soup.find('title').get_text() if soup.find('title') else None,
            "description": soup.find('meta', attrs={'name': 'description'})['content'] if soup.find('meta', attrs={'name': 'description'}) else None,
            "keywords": soup.find('meta', attrs={'name': 'keywords'})['content'] if soup.find('meta', attrs={'name': 'keywords'}) else None,
            "robots": soup.find('meta', attrs={'name': 'robots'})['content'] if soup.find('meta', attrs={'name': 'robots'}) else None,
            "canonical": soup.find('link', attrs={'rel': 'canonical'})['href'] if soup.find('link', attrs={'rel': 'canonical'}) else None
        }
        
        return meta_tags
    except Exception as e:
        raise Exception(f"Meta tag analysis failed: {str(e)}")

async def generate_optimized_meta_tags_content(current_meta: dict, keywords: str) -> dict:
    """Generate optimized meta tags based on current ones and keywords"""
    try:
        # Create meta content generation request
        meta_request = ContentGenerationRequest(
            content_type="meta-content",
            keywords=keywords,
            tone="professional",
            length="short",
            audience="business"
        )
        
        # Generate optimized meta content
        meta_response = generate_content(meta_request)
        
        # Parse the generated content to extract meta tags
        content = meta_response.content
        
        # Extract title and description from generated content
        title_match = re.search(r'Title[:\s]+(.+)', content, re.IGNORECASE)
        description_match = re.search(r'Description[:\s]+(.+)', content, re.IGNORECASE)
        
        optimized_meta = {
            "title": title_match.group(1).strip() if title_match else current_meta.get("title", ""),
            "description": description_match.group(1).strip() if description_match else current_meta.get("description", ""),
            "keywords": keywords,
            "robots": "index, follow",
            "canonical": current_meta.get("canonical", "")
        }
        
        return optimized_meta
    except Exception as e:
        raise Exception(f"Meta tag generation failed: {str(e)}")

def generate_meta_implementation_code(optimized_meta: dict) -> str:
    """Generate HTML code for implementing optimized meta tags"""
    html_code = f"""<!-- Optimized Meta Tags -->
<title>{optimized_meta.get('title', '')}</title>
<meta name="description" content="{optimized_meta.get('description', '')}">
<meta name="keywords" content="{optimized_meta.get('keywords', '')}">
<meta name="robots" content="{optimized_meta.get('robots', 'index, follow')}">
<link rel="canonical" href="{optimized_meta.get('canonical', '')}">

<!-- Open Graph Meta Tags -->
<meta property="og:title" content="{optimized_meta.get('title', '')}">
<meta property="og:description" content="{optimized_meta.get('description', '')}">
<meta property="og:type" content="website">

<!-- Twitter Card Meta Tags -->
<meta name="twitter:card" content="summary">
<meta name="twitter:title" content="{optimized_meta.get('title', '')}">
<meta name="twitter:description" content="{optimized_meta.get('description', '')}">"""
    
    return html_code

async def run_comprehensive_audit(url: str) -> dict:
    """Run a comprehensive SEO audit"""
    try:
        request = SEOAuditRequest(url=url)
        audit_result = seo_audit(request)
        return audit_result.dict()
    except Exception as e:
        raise Exception(f"Comprehensive audit failed: {str(e)}")

async def generate_implementation_plan(audit_result: dict, priority_level: str) -> dict:
    """Generate a comprehensive implementation plan"""
    seo_score = audit_result.get('seo_score', 0)
    content_score = audit_result.get('content_quality_score', 0)
    technical_score = audit_result.get('technical_score', 0)
    
    plan = {
        "overview": {
            "current_seo_score": seo_score,
            "target_seo_score": min(95, seo_score + 20),
            "priority_level": priority_level,
            "estimated_timeline": "4-6 weeks"
        },
        "critical_actions": [],
        "important_actions": [],
        "nice_to_have": []
    }
    
    # Critical actions (SEO score < 50)
    if seo_score < 50:
        plan["critical_actions"].extend([
            "Fix broken links and images",
            "Add missing meta descriptions",
            "Improve content quality and length",
            "Fix heading structure issues"
        ])
    
    # Important actions (SEO score 50-75)
    if 50 <= seo_score < 75:
        plan["important_actions"].extend([
            "Optimize meta tags",
            "Improve content structure",
            "Add internal linking",
            "Enhance user experience"
        ])
    
    # Nice to have (SEO score > 75)
    if seo_score >= 75:
        plan["nice_to_have"].extend([
            "Add structured data",
            "Optimize for featured snippets",
            "Enhance social media presence",
            "Create content calendar"
        ])
    
    return plan

def generate_implementation_timeline(plan: dict) -> dict:
    """Generate timeline for implementation plan"""
    return {
        "week_1": {
            "focus": "Critical fixes and content structure",
            "tasks": plan.get("critical_actions", [])[:3]
        },
        "week_2": {
            "focus": "Content optimization and meta tags",
            "tasks": plan.get("important_actions", [])[:3]
        },
        "week_3": {
            "focus": "Technical improvements and monitoring",
            "tasks": plan.get("important_actions", [])[3:6] if len(plan.get("important_actions", [])) > 3 else []
        },
        "week_4": {
            "focus": "Enhancement and fine-tuning",
            "tasks": plan.get("nice_to_have", [])[:3]
        }
    }

def generate_resource_requirements(plan: dict) -> dict:
    """Generate resource requirements for implementation"""
    return {
        "content_writer": "Required for content optimization",
        "web_developer": "Required for technical fixes",
        "seo_specialist": "Recommended for advanced optimization",
        "designer": "Optional for visual improvements",
        "estimated_hours": "20-40 hours depending on scope"
    }

async def analyze_html_structure(url: str) -> dict:
    """Analyze current HTML structure of a website"""
    try:
        response = requests.get(url, timeout=10)
        soup = BeautifulSoup(response.content, 'html.parser')
        
        analysis = {
            "headings": {
                "h1_count": len(soup.find_all('h1')),
                "h2_count": len(soup.find_all('h2')),
                "h3_count": len(soup.find_all('h3')),
                "heading_structure": "Good" if len(soup.find_all('h1')) == 1 else "Needs improvement"
            },
            "content": {
                "word_count": len(soup.get_text().split()),
                "paragraph_count": len(soup.find_all('p')),
                "content_quality": "Good" if len(soup.get_text().split()) > 300 else "Needs improvement"
            },
            "images": {
                "total_images": len(soup.find_all('img')),
                "images_with_alt": len([img for img in soup.find_all('img') if img.get('alt')]),
                "alt_text_coverage": "Good" if len([img for img in soup.find_all('img') if img.get('alt')]) / max(len(soup.find_all('img')), 1) > 0.8 else "Needs improvement"
            },
            "links": {
                "internal_links": len([a for a in soup.find_all('a') if a.get('href') and a.get('href').startswith('/')]),
                "external_links": len([a for a in soup.find_all('a') if a.get('href') and a.get('href').startswith('http')]),
                "link_quality": "Good" if len([a for a in soup.find_all('a') if a.get('href')]) > 5 else "Needs improvement"
            }
        }
        
        return analysis
    except Exception as e:
        raise Exception(f"HTML structure analysis failed: {str(e)}")

async def generate_optimized_html(current_html: dict, optimization_focus: str) -> dict:
    """Generate optimized HTML based on current structure"""
    recommendations = []
    
    if optimization_focus == "content":
        if current_html["content"]["word_count"] < 300:
            recommendations.append("Increase content length to at least 500 words")
        if current_html["headings"]["h1_count"] != 1:
            recommendations.append("Ensure only one H1 heading per page")
        if current_html["headings"]["h2_count"] < 2:
            recommendations.append("Add more H2 headings for better structure")
    
    elif optimization_focus == "structure":
        if current_html["headings"]["heading_structure"] != "Good":
            recommendations.append("Improve heading hierarchy (H1 > H2 > H3)")
        if current_html["links"]["link_quality"] != "Good":
            recommendations.append("Add more internal and external links")
    
    elif optimization_focus == "accessibility":
        if current_html["images"]["alt_text_coverage"] != "Good":
            recommendations.append("Add alt text to all images")
        recommendations.append("Ensure proper ARIA labels")
        recommendations.append("Improve color contrast")
    
    return {
        "current_analysis": current_html,
        "optimization_recommendations": recommendations,
        "html_code_examples": generate_html_code_examples(recommendations)
    }

def generate_html_code_examples(recommendations: list) -> dict:
    """Generate HTML code examples for recommendations"""
    examples = {}
    
    for rec in recommendations:
        if "H1 heading" in rec:
            examples["heading_structure"] = """<!-- Good heading structure -->
<h1>Main Page Title</h1>
<h2>Section Heading</h2>
<h3>Subsection Heading</h3>"""
        
        elif "alt text" in rec:
            examples["image_alt"] = """<!-- Image with proper alt text -->
<img src="image.jpg" alt="Descriptive text about the image" />"""
        
        elif "internal links" in rec:
            examples["internal_links"] = """<!-- Internal linking -->
<a href="/services">Our Services</a>
<a href="/about">About Us</a>
<a href="/contact">Contact Us</a>"""
    
    return examples

def generate_changes_summary(current_html: dict, optimized_html: dict) -> dict:
    """Generate summary of changes needed"""
    return {
        "priority_changes": optimized_html.get("optimization_recommendations", [])[:3],
        "estimated_impact": "High" if len(optimized_html.get("optimization_recommendations", [])) > 5 else "Medium",
        "implementation_time": "2-4 hours",
        "seo_improvement": "10-20 points"
    }